#!/usr/bin/env python3
"""Generate Rossi Mission Project Slide Deck V2 — Palette System Validated."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -- Color Palette --
DARK_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
MID_NAVY = RGBColor(0x2D, 0x2D, 0x4A)
ACCENT_BLUE = RGBColor(0x3A, 0x6B, 0x9F)
ACCENT_TEAL = RGBColor(0x2A, 0x9D, 0x8F)
ACCENT_ORANGE = RGBColor(0xE7, 0x6F, 0x51)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM_BG = RGBColor(0xFA, 0xF8, 0xF5)
EVIDENCE_COLOR = RGBColor(0x66, 0x88, 0x66)


def add_bg_rect(slide, color=DARK_NAVY):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, italic=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=DARK_TEXT, line_spacing=1.3, font_name='Calibri',
                   alignment=PP_ALIGN.LEFT):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, bold, italic) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1) + 4)
    return txBox


def add_evidence(slide, left, top, width, text):
    """Add an evidence marker in Yuty style."""
    return add_text_box(slide, left, top, width, 0.3, text,
                        font_size=9, italic=True, color=EVIDENCE_COLOR)


def add_accent_bar(slide, left, top, width, height, color=ACCENT_TEAL):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def section_slide(title, subtitle=None):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_rect(slide, DARK_NAVY)
    add_accent_bar(slide, 1.5, 3.2, 2.5, 0.06, ACCENT_TEAL)
    add_text_box(slide, 1.5, 3.4, 10, 1.2, title,
                 font_size=36, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, 1.5, 4.6, 10, 0.8, subtitle,
                     font_size=18, italic=True, color=RGBColor(0xAA, 0xAA, 0xCC))
    return slide


def content_slide(title, body_lines=None, accent_color=ACCENT_TEAL):
    """Create a standard content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_rect(slide, WARM_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.08, accent_color)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, title,
                 font_size=28, bold=True, color=DARK_NAVY)
    add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, accent_color)
    if body_lines:
        add_multi_text(slide, 0.8, 1.5, 11.5, 5.5, body_lines,
                       font_size=17, color=DARK_TEXT)
    return slide


def two_column_slide(title, left_title, left_lines, right_title, right_lines,
                     accent_color=ACCENT_TEAL):
    """Create a two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WARM_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.08, accent_color)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, title,
                 font_size=28, bold=True, color=DARK_NAVY)
    add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, accent_color)

    # Left column
    add_text_box(slide, 0.8, 1.5, 5.5, 0.5, left_title,
                 font_size=20, bold=True, color=ACCENT_BLUE)
    add_multi_text(slide, 0.8, 2.1, 5.5, 5.0, left_lines,
                   font_size=15, color=DARK_TEXT)

    # Divider
    add_accent_bar(slide, 6.6, 1.5, 0.03, 5.0, MED_GRAY)

    # Right column
    add_text_box(slide, 7.0, 1.5, 5.5, 0.5, right_title,
                 font_size=20, bold=True, color=ACCENT_BLUE)
    add_multi_text(slide, 7.0, 2.1, 5.5, 5.0, right_lines,
                   font_size=15, color=DARK_TEXT)
    return slide


def stat_box(slide, left, top, width, height, number, label,
             num_color=ACCENT_TEAL, bg_color=WHITE):
    """Add a stat callout box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = num_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = MED_GRAY
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER
    return shape


# ================================================================
# SLIDE 1: TITLE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, DARK_NAVY)

add_text_box(slide, 1.5, 1.5, 10, 1.0,
             'ROSSI MISSION PROJECT',
             font_size=42, bold=True, color=WHITE)
add_accent_bar(slide, 1.5, 2.7, 3.0, 0.06, ACCENT_TEAL)
add_text_box(slide, 1.5, 3.0, 10, 0.8,
             'Product Strategy & Discovery Plan',
             font_size=24, bold=False, color=RGBColor(0xCC, 0xCC, 0xDD))
add_text_box(slide, 1.5, 3.8, 10, 0.6,
             'A product-craft approach to validating and scaling a cultural platform for graffiti artists',
             font_size=16, italic=True, color=RGBColor(0x99, 0x99, 0xBB))

add_text_box(slide, 1.5, 5.2, 5, 0.4,
             'V2 \u2014 Palette System Validated',
             font_size=14, bold=True, color=ACCENT_TEAL)
add_text_box(slide, 1.5, 5.6, 5, 0.4,
             'Prepared for Rossi Leadership Team',
             font_size=14, color=RGBColor(0x88, 0x88, 0xAA))
add_text_box(slide, 1.5, 5.9, 5, 0.4,
             'February 2026  |  Product Management Practice',
             font_size=12, color=RGBColor(0x77, 0x77, 0x99))
add_text_box(slide, 9.5, 6.5, 3, 0.4,
             'CONFIDENTIAL',
             font_size=10, bold=True, color=ACCENT_RED,
             alignment=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 2: AGENDA
# ================================================================
slide = content_slide('Agenda')
items = [
    ('1.  The Customer Problem \u2014 Who we serve and why it matters', False, False),
    ('2.  The Market \u2014 Who Rossi competes with (named competitors + data)', False, False),
    ('3.  What We Know \u2014 Evidence from market research and operations', False, False),
    ('4.  How Top Artists Sell Today \u2014 Real revenue models, real numbers', False, False),
    ('5.  What We Believe \u2014 Hypotheses that need testing', False, False),
    ('6.  The Product \u2014 What Rossi actually sells', False, False),
    ('7.  Product Strategy \u2014 Year 1-3 approach with partner targets', False, False),
    ('8.  Product Risks \u2014 And how we surface them early', False, False),
    ('9.  The 90-Day Discovery Plan \u2014 $10K to transform plan into evidence', False, False),
    ('10. Recommendation \u2014 Conditional Go', True, False),
]
add_multi_text(slide, 0.8, 1.5, 11.5, 5.5, items,
               font_size=18, color=DARK_TEXT, line_spacing=1.4)

# ================================================================
# SLIDE 3: SECTION — THE CUSTOMER PROBLEM
# ================================================================
section_slide('The Customer Problem',
              'Before we talk about the product, we need to talk about the customer.')

# ================================================================
# SLIDE 4: THE ARTIST'S PROBLEM
# ================================================================
slide = content_slide('The Artist\'s Problem')
lines = [
    ('Graffiti artists face a career ceiling that other artists do not.', True, False),
    ('', False, False),
    ('Traditional galleries treat street art as illegitimate \u2014 no gallery path, no museum trajectory.', False, False),
    ('', False, False),
    ('Streetwear brands pay flat fees for designs. Artists are vendors, not partners.', False, False),
    ('  KAWS x Uniqlo, Futura x Kenzo, Stash x Nike \u2014 licensing fees $10K-$500K, royalties 5-15%.', False, True),
    ('  The artist creates the value. The brand captures most of it.', False, True),
    ('', False, False),
    ('Online platforms offer distribution but no curation, career development, or community.', False, False),
    ('  Etsy: 10,000+ graffiti listings, 6.5% fee \u2014 but no authentication, race to bottom.', False, True),
    ('', False, False),
    ('Result: Artists with signature styles and cultural credibility who cannot convert that into sustainable careers.', True, True),
]
add_multi_text(slide, 0.8, 1.5, 11.5, 5.5, lines,
               font_size=16, color=DARK_TEXT, line_spacing=1.1)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Marketplace Research \u2014 Brand Collaborations + Online Platforms]')

# ================================================================
# SLIDE 5: THE PROOF — IT CAN WORK
# ================================================================
slide = content_slide('The Proof: It Can Work')
lines = [
    ('Barry McGee: Mission District tagger \u2192 Venice Biennale (2001)', False, False),
    ('Chito: Custom jacket reworks \u2192 Givenchy collection (2020)', False, False),
    ('KAWS: Brooklyn tagger \u2192 $15M auction record, 4.8M Instagram followers', False, False),
    ('Shepard Fairey: Sticker campaign \u2192 OBEY (full brand ownership, 2.1M followers)', False, False),
    ('', False, False),
    ('They all had the same formula:', True, False),
    ('    Signature style + Champion + Documentation = Career launch', False, True),
    ('', False, False),
    ('But they got there through luck and individual connections, not through any system.', False, False),
    ('', False, False),
    ('The problem statement:', True, False),
    ('There is no institution that systematically discovers, develops, and launches graffiti artist careers while keeping artists as economic partners.', True, True),
]
add_multi_text(slide, 0.8, 1.5, 11.5, 5.5, lines,
               font_size=16, color=DARK_TEXT, line_spacing=1.1)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Marketplace Research \u2014 Tier 1 Artists; BCG Assessment]')

# ================================================================
# SLIDE 6: SECTION — THE MARKET
# ================================================================
section_slide('The Market',
              'Who Rossi competes with \u2014 named competitors, real data.')

# ================================================================
# SLIDE 7: COMPETITIVE LANDSCAPE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'The Competitive Landscape: Named Competitors',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_BLUE)

# Competitor boxes
competitors = [
    ('GraffitiStreet.com', '100+ artists\n40-50% commission\nGlobal shipping\nOnline gallery only',
     'Banksy, Invader,\nOs Gemeos, Spidertag', ACCENT_BLUE),
    ('DirtyPilot.com', 'Curated prints\nAuthentication certs\nLimited editions\nNo physical space',
     'Cope2, Lady Pink', ACCENT_TEAL),
    ('Etsy', '10,000+ listings\n6.5% fee\nNo curation\nAuth problems',
     'Mass market\nUnauthorized reprods', ACCENT_ORANGE),
    ('StockX', 'Authenticated sales\n9.5% seller fee\nAuction model\nCollectibles focus',
     'KAWS, Futura,\nCope2', ACCENT_RED),
]

for i, (name, details, artists, color) in enumerate(competitors):
    x = 0.5 + i * 3.15
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5),
        Inches(2.9), Inches(4.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = ''
    p3 = tf.add_paragraph()
    p3.text = details
    p3.font.size = Pt(11)
    p3.font.color.rgb = DARK_TEXT
    p3.font.name = 'Calibri'
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = ''
    p5 = tf.add_paragraph()
    p5.text = 'Key Artists:'
    p5.font.size = Pt(10)
    p5.font.bold = True
    p5.font.color.rgb = MED_GRAY
    p5.alignment = PP_ALIGN.CENTER

    p6 = tf.add_paragraph()
    p6.text = artists
    p6.font.size = Pt(10)
    p6.font.color.rgb = MED_GRAY
    p6.alignment = PP_ALIGN.CENTER

# Rossi advantage callout
add_multi_text(slide, 0.8, 5.8, 11.5, 1.5, [
    ('Rossi\'s differentiator: No competitor combines curation + physical gallery + artist development + 50/50 economics.', True, False),
    ('GraffitiStreet has scale. Etsy has reach. Neither builds artist careers.', False, True),
], font_size=14, color=DARK_TEXT)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Marketplace Research \u2014 Competitive Landscape]')

# ================================================================
# SLIDE 8: HOW TOP ARTISTS SELL TODAY
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_TEAL)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'How Top Artists Actually Sell Today',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_TEAL)

# Artist data table - headers
headers = ['Artist', 'Followers', 'Channels', 'Price Range', 'Model']
col_widths = [2.0, 1.5, 3.5, 2.2, 3.0]
x_positions = [0.6]
for w in col_widths[:-1]:
    x_positions.append(x_positions[-1] + w)

for j, header in enumerate(headers):
    add_text_box(slide, x_positions[j], 1.5, col_widths[j], 0.4, header,
                 font_size=12, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# Header background
hdr_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.45),
    Inches(12.3), Inches(0.45)
)
hdr_bg.fill.solid()
hdr_bg.fill.fore_color.rgb = MID_NAVY
hdr_bg.line.fill.background()
# Move header bg to back
slide.shapes._spTree.remove(hdr_bg._element)
slide.shapes._spTree.insert(2, hdr_bg._element)

# Re-add headers on top
for j, header in enumerate(headers):
    add_text_box(slide, x_positions[j], 1.5, col_widths[j], 0.4, header,
                 font_size=12, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# Artist rows
artists_data = [
    ('KAWS', '4.8M', 'DDTStore, Brooklyn Museum, StockX', '$200-$50K', 'Toys + drops + brand collabs'),
    ('Shepard Fairey', '2.1M', 'OBEYClothing.com, galleries', '$30-$100K', 'Full brand ownership'),
    ('Cope2', '400K', 'Cope2-Merch.com, DirtyPilot', '$200-$25K', 'D2C + gallery + prints'),
    ('Lady Pink', '200K', 'LadyPinkNYC.com, Wynwood', '$50-$50K', 'D2C + gallery + brand collabs'),
    ('Futura 2000', '500K', 'Wynwood Shop, StockX, Kenzo', '$100-$200K', 'Fashion collabs + gallery'),
]

for i, (artist, followers, channels, price, model) in enumerate(artists_data):
    y = 2.0 + i * 0.55
    row_data = [artist, followers, channels, price, model]
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT]
    for j, val in enumerate(row_data):
        add_text_box(slide, x_positions[j], y, col_widths[j], 0.5, val,
                     font_size=11, color=DARK_TEXT, alignment=aligns[j])
    # Alternating row bg
    if i % 2 == 0:
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y - 0.02),
            Inches(12.3), Inches(0.5)
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)
        row_bg.line.fill.background()
        slide.shapes._spTree.remove(row_bg._element)
        slide.shapes._spTree.insert(2, row_bg._element)

add_multi_text(slide, 0.8, 4.9, 11.5, 1.8, [
    ('Key pattern: Every successful artist uses 3-5 channels simultaneously.', True, False),
    ('Most use hybrid model: D2C website + Instagram drops + gallery + brand collabs.', False, False),
    ('Rossi\'s multi-channel approach is validated by how these artists already operate.', False, True),
], font_size=14, color=DARK_TEXT, line_spacing=1.2)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Marketplace Research \u2014 Tier 1 & 2 Artist Profiles, Instagram Data]')

# ================================================================
# SLIDE 9: REVENUE MODELS ARE PROVEN
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_ORANGE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Six Revenue Models \u2014 All Proven in the Wild',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_ORANGE)

models = [
    ('D2C Shops', 'Cope2-Merch.com\nSaberOne.com\nLadyPinkNYC.com', '70-90%\nmargin', ACCENT_TEAL),
    ('Gallery Rep', 'Retna, Mear One\nOs Gemeos', '50-60%\nto artist', ACCENT_BLUE),
    ('Print Editions', 'DirtyPilot\nEdition Copenhagen', '50/50\nsplit', ACCENT_ORANGE),
    ('Marketplaces', 'Etsy, Redbubble\nStockX', '80-94%\nto artist', ACCENT_GREEN),
    ('Brand Collabs', 'KAWS x Dior\nFutura x Kenzo', '$10K-$500K\nlicensing', ACCENT_RED),
    ('Collectives', 'Alberta Street\nSide Rail', '70-100%\nto artist', RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (model_name, examples, economics, color) in enumerate(models):
    x = 0.4 + (i % 3) * 4.2
    y = 1.5 + (i // 3) * 2.8
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
        Inches(3.9), Inches(2.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = model_name
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = examples
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_TEXT
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = ''
    p4 = tf.add_paragraph()
    p4.text = economics
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = color
    p4.alignment = PP_ALIGN.CENTER

add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Marketplace Research \u2014 Revenue Models Observed]')

# ================================================================
# SLIDE 10: SECTION — WHAT WE KNOW
# ================================================================
section_slide('What We Know', 'Evidence, not opinions.')

# ================================================================
# SLIDE 11: MARKET EVIDENCE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'The Market Is Real',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_BLUE)

stat_box(slide, 0.8, 1.6, 2.7, 1.6, '$6.8B', 'Global Street Art Market\n8.2% CAGR', ACCENT_TEAL)
stat_box(slide, 3.8, 1.6, 2.7, 1.6, '50+', 'Active Graffiti Artists\nWith Merch Operations', ACCENT_BLUE)
stat_box(slide, 6.8, 1.6, 2.7, 1.6, '6', 'Proven Revenue Models\nAll Active in Market', ACCENT_ORANGE)
stat_box(slide, 9.8, 1.6, 2.7, 1.6, '2-4%', 'Valencia St Vacancy\n(vs 6-8% citywide)', ACCENT_GREEN)

add_multi_text(slide, 0.8, 3.5, 11.5, 4.0, [
    ('Creative Growth Art Center validates the exact model at scale:', True, False),
    ('  Nonprofit gallery  |  50/50 artist split  |  140 artists  |  $3.26M revenue  |  50 years', False, False),
    ('  Form 990 data (EIN 23-7319028): Art sales ranged $499K\u2013$1.06M over 5 years (volatile).', False, False),
    ('  Grants provided 30\u201341% of revenue every year \u2014 the stabilizer. Personnel: 57% of costs.', False, False),
    ('', False, False),
    ('What Rossi already has:', True, False),
    ('  80+ roster artists  |  251 SKUs  |  791 Valencia St  |  Operating Shopify store', False, False),
    ('', False, False),
    ('"Rossi is the graffiti/street art version of Creative Growth, serving an underserved', False, True),
    ('  artist community in SF\u2019s Mission District." \u2014 Infinitely more credible than Supreme or St\u00fcssy.', False, True),
], font_size=15, color=DARK_TEXT, line_spacing=1.1)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Creative Growth Form 990 (2020\u20132024); Marketplace Research (commission rates)]')

# ================================================================
# SLIDE 12: THE CRITICAL GAP
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_RED)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'The Critical Gap',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_RED)

# Big callout
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8),
    Inches(10.3), Inches(2.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xF0)
shape.line.color.rgb = ACCENT_RED
shape.line.width = Pt(2)

tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'We do not know Rossi\'s current revenue.'
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = ACCENT_RED
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = ('The business plan projects $780K Year 1, but no one has pulled the POS data.\n'
           'Modeled estimate: $245K-$390K annually. The gap between $250K and $400K\n'
           'determines whether the Year 1 target is ambitious or unrealistic.')
p2.font.size = Pt(16)
p2.font.color.rgb = DARK_TEXT
p2.font.name = 'Calibri'
p2.alignment = PP_ALIGN.CENTER

add_multi_text(slide, 0.8, 4.5, 11.5, 2.5, [
    ('The Underwriter\'s verdict: "Analysis is not evidence. This is a world-class plan. It needs real data."', False, True),
    ('', False, False),
    ('Meanwhile: Cope2 runs Cope2-Merch.com with direct D2C sales. Lady Pink sells through LadyPinkNYC.com.', False, False),
    ('Artists are already proving online revenue works \u2014 Rossi just hasn\'t measured its own.', True, False),
], font_size=16, color=DARK_TEXT, line_spacing=1.2)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Underwriter Audit \u2014 82/100; Marketplace Research \u2014 D2C Shops]')

# ================================================================
# SLIDE 13: SECTION — HYPOTHESES
# ================================================================
section_slide('What We Believe',
              'Hypotheses are essential. But they must be tested before we build on them.')

# ================================================================
# SLIDE 14: KEY HYPOTHESES
# ================================================================
slide = two_column_slide(
    'Core Hypotheses \u2014 Ranked by Risk',
    'HIGH RISK (Unvalidated)',
    [
        ('Workshops discover talent (10% conversion)', True, False),
        ('  Zero workshops conducted', False, True),
        ('  Test: 1 pilot workshop, track 90 days', False, False),
        ('', False, False),
        ('Content drives art sales (GX1000 model)', True, False),
        ('  No attribution data exists', False, True),
        ('  Test: Tag links, track funnel 60 days', False, False),
        ('', False, False),
        ('Residencies launch careers (80% rate)', True, False),
        ('  Zero residencies completed', False, True),
        ('  Test: Run 1 residency, document outcomes', False, False),
    ],
    'MEDIUM RISK (Partially Validated)',
    [
        ('Scarcity works for emerging artists', True, False),
        ('  KAWS proved it (4.8M followers, $200-$50K range)', False, True),
        ('  But Rossi is not KAWS. Test: monitor sell-through 6 months', False, False),
        ('', False, False),
        ('Artists stay for 50/50 split', True, False),
        ('  Alberta St Gallery uses 70/30. Side Rail uses 100%.', False, True),
        ('  Rossi must prove 50/50 + development = better deal', False, False),
        ('', False, False),
        ('Partner galleries will join network', True, False),
        ('  Real targets identified (see Strategy slide)', False, True),
        ('  Test: Pitch 3 galleries, get 1 LOI', False, False),
    ],
    accent_color=ACCENT_ORANGE
)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: BCG Assessment \u2014 Risk Register; Marketplace Research \u2014 Commission Rates]')

# ================================================================
# SLIDE 15: SECTION — THE PRODUCT
# ================================================================
section_slide('The Product',
              'What Rossi actually sells \u2014 three products, one platform.')

# ================================================================
# SLIDE 16: THREE PRODUCTS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_TEAL)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Rossi Sells Three Products',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_TEAL)

# Product 1 box
shape1 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6),
    Inches(3.9), Inches(5.2)
)
shape1.fill.solid()
shape1.fill.fore_color.rgb = WHITE
shape1.line.color.rgb = ACCENT_TEAL
shape1.line.width = Pt(2)

tf1 = shape1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = 'CURATED ART'
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER
p2 = tf1.add_paragraph()
p2.text = 'Gallery + Online'
p2.font.size = Pt(14)
p2.font.color.rgb = MED_GRAY
p2.alignment = PP_ALIGN.CENTER
p3 = tf1.add_paragraph()
p3.text = ''
p4 = tf1.add_paragraph()
p4.text = '47% of Year 1 Revenue'
p4.font.size = Pt(16)
p4.font.bold = True
p4.font.color.rgb = DARK_TEXT
p4.alignment = PP_ALIGN.CENTER
p5 = tf1.add_paragraph()
p5.text = ''
p6 = tf1.add_paragraph()
p6.text = ('Originals, limited prints, merch\n250 SKU cap per location\n50/50 artist split\n'
           'Physical + Shopify')
p6.font.size = Pt(13)
p6.font.color.rgb = DARK_TEXT
p6.alignment = PP_ALIGN.CENTER
p7 = tf1.add_paragraph()
p7.text = ''
p8 = tf1.add_paragraph()
p8.text = 'Benchmark: GraffitiStreet (40-50%)'
p8.font.size = Pt(11)
p8.font.color.rgb = EVIDENCE_COLOR
p8.alignment = PP_ALIGN.CENTER
p9 = tf1.add_paragraph()
p9.text = ''
p10 = tf1.add_paragraph()
p10.text = 'PMF Status: VALIDATED'
p10.font.size = Pt(13)
p10.font.bold = True
p10.font.color.rgb = ACCENT_GREEN
p10.alignment = PP_ALIGN.CENTER
p11 = tf1.add_paragraph()
p11.text = '(Product exists, customers buy)'
p11.font.size = Pt(11)
p11.font.color.rgb = MED_GRAY
p11.alignment = PP_ALIGN.CENTER

# Product 2 box
shape2 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.6),
    Inches(3.9), Inches(5.2)
)
shape2.fill.solid()
shape2.fill.fore_color.rgb = WHITE
shape2.line.color.rgb = ACCENT_BLUE
shape2.line.width = Pt(2)

tf2 = shape2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = 'EXPERIENCES'
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER
p2 = tf2.add_paragraph()
p2.text = 'Workshops + Events + Residencies'
p2.font.size = Pt(14)
p2.font.color.rgb = MED_GRAY
p2.alignment = PP_ALIGN.CENTER
p3 = tf2.add_paragraph()
p3.text = ''
p4 = tf2.add_paragraph()
p4.text = '38% of Year 1 Revenue'
p4.font.size = Pt(16)
p4.font.bold = True
p4.font.color.rgb = DARK_TEXT
p4.alignment = PP_ALIGN.CENTER
p5 = tf2.add_paragraph()
p5.text = ''
p6 = tf2.add_paragraph()
p6.text = ('Public + corporate workshops\nMonthly gallery events\nSponsored residencies\n'
           'Talent discovery pipeline')
p6.font.size = Pt(13)
p6.font.color.rgb = DARK_TEXT
p6.alignment = PP_ALIGN.CENTER
p7 = tf2.add_paragraph()
p7.text = ''
p8 = tf2.add_paragraph()
p8.text = 'Benchmark: Creative Growth ($3.26M)'
p8.font.size = Pt(11)
p8.font.color.rgb = EVIDENCE_COLOR
p8.alignment = PP_ALIGN.CENTER
p9 = tf2.add_paragraph()
p9.text = ''
p10 = tf2.add_paragraph()
p10.text = 'PMF Status: UNVALIDATED'
p10.font.size = Pt(13)
p10.font.bold = True
p10.font.color.rgb = ACCENT_RED
p10.alignment = PP_ALIGN.CENTER
p11 = tf2.add_paragraph()
p11.text = '(Zero workshops or residencies run)'
p11.font.size = Pt(11)
p11.font.color.rgb = MED_GRAY
p11.alignment = PP_ALIGN.CENTER

# Product 3 box
shape3 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.6),
    Inches(3.9), Inches(5.2)
)
shape3.fill.solid()
shape3.fill.fore_color.rgb = WHITE
shape3.line.color.rgb = ACCENT_ORANGE
shape3.line.width = Pt(2)

tf3 = shape3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = 'BRAND PARTNERSHIPS'
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER
p2 = tf3.add_paragraph()
p2.text = 'Collaborations + Licensing'
p2.font.size = Pt(14)
p2.font.color.rgb = MED_GRAY
p2.alignment = PP_ALIGN.CENTER
p3 = tf3.add_paragraph()
p3.text = ''
p4 = tf3.add_paragraph()
p4.text = '15% of Year 1 Revenue'
p4.font.size = Pt(16)
p4.font.bold = True
p4.font.color.rgb = DARK_TEXT
p4.alignment = PP_ALIGN.CENTER
p5 = tf3.add_paragraph()
p5.text = ''
p6 = tf3.add_paragraph()
p6.text = ('3-5 per year (scarcity cap)\n70% artist / 30% Rossi\nArtist veto power\n'
           'Cultural fit required')
p6.font.size = Pt(13)
p6.font.color.rgb = DARK_TEXT
p6.alignment = PP_ALIGN.CENTER
p7 = tf3.add_paragraph()
p7.text = ''
p8 = tf3.add_paragraph()
p8.text = 'Benchmark: KAWS x Dior, Futura x Kenzo'
p8.font.size = Pt(11)
p8.font.color.rgb = EVIDENCE_COLOR
p8.alignment = PP_ALIGN.CENTER
p9 = tf3.add_paragraph()
p9.text = 'Licensing: $10K-$500K + 5-15% royalties'
p9.font.size = Pt(10)
p9.font.color.rgb = EVIDENCE_COLOR
p9.alignment = PP_ALIGN.CENTER
p10 = tf3.add_paragraph()
p10.text = ''
p11 = tf3.add_paragraph()
p11.text = 'PMF Status: UNVALIDATED'
p11.font.size = Pt(13)
p11.font.bold = True
p11.font.color.rgb = ACCENT_RED
p11.alignment = PP_ALIGN.CENTER
p12 = tf3.add_paragraph()
p12.text = '(No collaborations completed)'
p12.font.size = Pt(11)
p12.font.color.rgb = MED_GRAY
p12.alignment = PP_ALIGN.CENTER

# ================================================================
# SLIDE 16B: ARTIST ECONOMICS REALITY (PARETO)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_ORANGE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Artist Economics: The Honest Math',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_ORANGE)

# Three tier boxes
tier_data = [
    ('TOP 10%', '8 Artists', '$210\u2013245K gross (60\u201370%)',
     '$13K\u2013$15K net/artist', 'Revenue anchors.\nFund the operation.',
     ACCENT_GREEN),
    ('MIDDLE 30%', '24 Artists', '$70\u201388K gross (20\u201325%)',
     '$1,450\u2013$1,830 net/artist', 'Growth tier.\nPipeline is working.',
     ACCENT_BLUE),
    ('EMERGING 60%', '48 Artists', '$18\u201370K gross (5\u201320%)',
     '$185\u2013$730 net/artist', 'Development tier.\nBuilding toward mid.',
     ACCENT_ORANGE),
]

for i, (tier, count, gross, net, role, color) in enumerate(tier_data):
    x = 0.5 + i * 4.2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5),
        Inches(3.8), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = tier
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = count
    p2.font.size = Pt(14)
    p2.font.color.rgb = MED_GRAY
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = ''
    p4 = tf.add_paragraph()
    p4.text = gross
    p4.font.size = Pt(12)
    p4.font.color.rgb = DARK_TEXT
    p4.alignment = PP_ALIGN.CENTER

    p5 = tf.add_paragraph()
    p5.text = net
    p5.font.size = Pt(16)
    p5.font.bold = True
    p5.font.color.rgb = color
    p5.alignment = PP_ALIGN.CENTER

    p6 = tf.add_paragraph()
    p6.text = ''
    p7 = tf.add_paragraph()
    p7.text = role
    p7.font.size = Pt(12)
    p7.font.color.rgb = MED_GRAY
    p7.font.name = 'Calibri'
    p7.alignment = PP_ALIGN.CENTER

add_multi_text(slide, 0.8, 5.4, 11.5, 1.8, [
    ('This is not a bug \u2014 this is how art markets work. Creative Growth shows the same pattern.', True, False),
    ('The pipeline\u2019s job: move artists UP the distribution. Workshops \u2192 visibility \u2192 sales \u2192 mid-tier \u2192 top.', False, False),
    ('The metric that matters: are artists moving up the curve year over year?', False, True),
], font_size=14, color=DARK_TEXT, line_spacing=1.15)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Validated Comparables \u2014 Pareto Distribution; Creative Growth 50-year pattern]')

# ================================================================
# SLIDE 17: BRAND ARCHITECTURE
# ================================================================
slide = content_slide('Recommended Brand Architecture')
lines = [
    ('Rossi Mission Project \u2014 The master brand and cultural platform', True, False),
    ('', False, False),
    ('Rossi Gallery \u2014 Physical and online retail experience', False, False),
    ('    Speaks to: Collectors, design-forward consumers', False, True),
    ('    Benchmark: GraffitiStreet (online) + 33 1/3 Gallery (physical)', False, True),
    ('', False, False),
    ('Rossi Studio \u2014 Artist development programs (workshops, residencies)', False, False),
    ('    Speaks to: Grant funders, emerging artists, corporate partners', False, True),
    ('    Benchmark: Creative Growth Art Center ($3.26M, 140 artists)', False, True),
    ('', False, False),
    ('Rossi Editions \u2014 Limited-edition drops (prints, merch, collaborations)', False, False),
    ('    Speaks to: Streetwear consumers, brand partners', False, True),
    ('    Benchmark: DirtyPilot (curated prints) + KAWS DDT drops', False, True),
    ('', False, False),
    ('Rossi Stories \u2014 Content and documentation engine', False, False),
    ('    Speaks to: Cultural audience, drives discovery across all products', False, True),
]
add_multi_text(slide, 0.8, 1.5, 11.5, 5.5, lines,
               font_size=15, color=DARK_TEXT, line_spacing=1.05)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Marketplace Research \u2014 Platform Models; BCG Assessment \u2014 Brand Architecture]')

# ================================================================
# SLIDE 18: SECTION — PRODUCT STRATEGY
# ================================================================
section_slide('Product Strategy',
              'Year 1 is about learning, not growth.')

# ================================================================
# SLIDE 19: YEAR 1 STRATEGY
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_TEAL)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Year 1: Validate the Core (SF Only)',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_TEAL)

stat_box(slide, 0.8, 1.5, 3.0, 1.3, '$745K', 'Year 1 Revenue Target', ACCENT_TEAL)
stat_box(slide, 4.0, 1.5, 3.0, 1.3, '1 Location', 'SF Flagship Only', ACCENT_BLUE)
stat_box(slide, 7.2, 1.5, 3.0, 1.3, '80+', 'Artists on Roster', ACCENT_ORANGE)

add_multi_text(slide, 0.8, 3.3, 11.5, 4.0, [
    ('Revenue Mix (Discovery-Informed):', True, False),
    ('  Retail (Gallery + Online): $405K (54%)  \u2014  Core product, known to work', False, False),
    ('  Workshops + Events: $160K (21%)  \u2014  Must validate with pilot workshops', False, False),
    ('  Collaborations: $70K (10%)  \u2014  Build pipeline, close 2-3 deals', False, False),
    ('  Grants: $110K (15%)  \u2014  Apply for $300K+ to net $110K (25\u201350% success rates)', False, False),
    ('', False, False),
    ('Grant reality check: Expected value from $215K in applications = $70K.', True, False),
    ('Rossi must apply aggressively (SFAC, CAC, Zellerbach, Haas, SF Foundation) and hire grant writer.', False, True),
    ('', False, False),
    ('Year 1 is NOT about hitting $745K. It is about learning whether the three products work.', True, True),
], font_size=14, color=DARK_TEXT, line_spacing=1.1)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Validated Comparables \u2014 Grant Landscape; Underwriter Audit \u2014 Revenue Model]')

# ================================================================
# SLIDE 20: PARTNER GALLERY TARGETS
# ================================================================
slide = two_column_slide(
    'Partner Expansion: Real Targets, Real Addresses',
    'Year 2 Targets (Validate First)',
    [
        ('Oakland: Gray Loft Gallery', True, False),
        ('  489 25th St, Jingletown', False, False),
        ('  Artist collective, First Friday walks', False, True),
        ('', False, False),
        ('Portland: Alberta Street Gallery', True, False),
        ('  1829 NE Alberta St', False, False),
        ('  30 members, 70/30 split, cooperative model', False, True),
        ('', False, False),
        ('Seattle: Side Rail Collective', True, False),
        ('  5511 1/2 Airport Way S, Georgetown', False, False),
        ('  No-commission model, membership fees', False, True),
    ],
    'Year 3 Targets (If Evidence Supports)',
    [
        ('LA: 01 Gallery', True, False),
        ('  Melrose \u2014 first to exhibit graffiti artists', False, False),
        ('  Mear One exhibited here pre-Banksy US debut', False, True),
        ('', False, False),
        ('LA: 33 1/3 Gallery', True, False),
        ('  Silverlake \u2014 hosted Banksy\'s first US show', False, False),
        ('', False, False),
        ('NYC: One Art Space', True, False),
        ('  Tribeca \u2014 curated graffiti shows', False, False),
        ('  Cope2, Shepard Fairey exhibited', False, True),
        ('', False, False),
        ('NYC: Bottleneck Gallery', True, False),
        ('  Limited edition prints, gig posters', False, False),
    ],
    accent_color=ACCENT_TEAL
)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Marketplace Research \u2014 Physical Locations; BCG Assessment \u2014 Expansion Plan]')

# ================================================================
# SLIDE 21: YEAR 1-3 TRAJECTORY
# ================================================================
slide = two_column_slide(
    'Year 1-3 Trajectory: Evidence-Based Growth',
    'Our Projections (Discovery-Based)',
    [
        ('Year 1: $745K (baseline)', True, False),
        ('  SF only. Validate all 3 products.', False, False),
        ('', False, False),
        ('Year 2: $1.0-1.1M (+35-45%)', True, False),
        ('  Add Oakland partner (Gray Loft).', False, False),
        ('  Scale what works.', False, False),
        ('', False, False),
        ('Year 3: $1.4-1.6M (+35-45%)', True, False),
        ('  Add LA + 1 city. Expand deliberately.', False, False),
        ('', False, False),
        ('Growth rate: exceptional but achievable.', False, True),
    ],
    'Original Plan (Concerns)',
    [
        ('Year 1: $780K', True, False),
        ('Year 2: $1.55M (+99%)', True, False),
        ('Year 3: $2.7M (+74%)', True, False),
        ('', False, False),
        ('99% Y2 growth is venture-scale.', False, False),
        ('Assumes new locations generate', False, False),
        ('$300-400K immediately (no ramp).', False, False),
        ('', False, False),
        ('Nonprofits grow at 20-25%/year.', False, False),
        ('Exceptional nonprofits: 35-45%.', False, False),
        ('The original plan exceeds both.', False, True),
    ],
    accent_color=ACCENT_BLUE
)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Underwriter Audit \u2014 Growth Analysis; BCG Assessment \u2014 Revenue Projections]')

# ================================================================
# SLIDE 22: SECTION — RISKS
# ================================================================
section_slide('Product Risks',
              'We do not predict the future. We design experiments that surface problems early.')

# ================================================================
# SLIDE 23: TOP RISKS
# ================================================================
slide = two_column_slide(
    'Top Product Risks',
    'High Impact Risks',
    [
        ('Brand Identity Confusion', True, False),
        ('  Is Rossi a brand, nonprofit, or incubator?', False, False),
        ('  Fix: Leadership alignment Week 1.', False, True),
        ('', False, False),
        ('Month 17 Cash Crisis ($10K balance)', True, False),
        ('  3 days of operating cash.', False, False),
        ('  Fix: Staged funding. Tranche 2 at Month 6.', False, True),
        ('', False, False),
        ('Workshop Pipeline Fails', True, False),
        ('  10% conversion is unvalidated.', False, False),
        ('  Fix: Pilot workshop Month 2. Redesign if <8 attend.', False, True),
    ],
    'Medium Impact Risks',
    [
        ('Personnel Costs Will Escalate', True, False),
        ('  Creative Growth: 57% on staff. Rossi: 28%.', False, False),
        ('  As you scale, expect 45\u201355%. Plan Year 2\u20133.', False, True),
        ('', False, False),
        ('Art Sales Volatility (2x Swings)', True, False),
        ('  CG art sales: $499K (2023) \u2192 $1.06M (2024).', False, False),
        ('  Fix: Build 3-month reserve. Don\u2019t assume steady.', False, True),
        ('', False, False),
        ('Top Artists Leave for D2C', True, False),
        ('  Cope2, Lady Pink keep 70\u201390% margin on own shops.', False, False),
        ('  Fix: Prove 50/50 + development > solo D2C.', False, True),
        ('', False, False),
        ('Rent Burden ($120K/yr = 16% of revenue)', True, False),
        ('  Structural disadvantage vs Oakland-based CG.', False, False),
        ('  Fix: Verify lease. Negotiate or plan lower-cost expansion.', False, True),
    ],
    accent_color=ACCENT_RED
)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Validated Comparables \u2014 CG Form 990; Marketplace Research \u2014 D2C Margins]')

# ================================================================
# SLIDE 24: SECTION — DISCOVERY PLAN
# ================================================================
section_slide('The 90-Day Discovery Plan',
              '$10,000 to transform a business plan into validated evidence.')

# ================================================================
# SLIDE 25: DISCOVERY OVERVIEW
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_GREEN)
add_text_box(slide, 0.8, 0.4, 11, 0.8, '90-Day Discovery Roadmap',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_GREEN)

# Phase boxes
for i, (phase, days, items, gate, color) in enumerate([
    ('Phase 1: Foundation', 'Days 1-30', [
        'Pull POS data (2 days)',
        'Hire photographer ($2K)',
        'Optimize Shopify ($3.5K)',
        'Launch first Artist Drop',
        'Survey 20 roster artists',
        'Benchmark: Cope2/Lady Pink D2C',
    ], 'Gate: Do we know our baseline?', ACCENT_TEAL),
    ('Phase 2: Validation', 'Days 31-60', [
        'Run pilot workshop ($500)',
        'Complete 3 artist case studies',
        'Execute 2nd & 3rd Artist Drops',
        'Pitch Alberta St + Gray Loft',
        'Submit grant applications',
        'Track online attribution data',
    ], 'Gate: Does the pipeline work?', ACCENT_BLUE),
    ('Phase 3: Package', 'Days 61-90', [
        'Update plan with real data',
        'Write Section 8 (financials)',
        'Reconcile all projections',
        'Prepare funding presentation',
        'Present to board + funders',
        'Competitive positioning deck',
    ], 'Gate: Does evidence support funding?', ACCENT_ORANGE),
]):
    x = 0.5 + i * 4.2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5),
        Inches(3.8), Inches(5.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = days
    p2.font.size = Pt(12)
    p2.font.color.rgb = MED_GRAY
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = ''

    for item in items:
        pi = tf.add_paragraph()
        pi.text = f'  {item}'
        pi.font.size = Pt(12)
        pi.font.color.rgb = DARK_TEXT
        pi.font.name = 'Calibri'

    pg = tf.add_paragraph()
    pg.text = ''
    pg2 = tf.add_paragraph()
    pg2.text = gate
    pg2.font.size = Pt(11)
    pg2.font.bold = True
    pg2.font.color.rgb = color
    pg2.font.name = 'Calibri'
    pg2.alignment = PP_ALIGN.CENTER

# ================================================================
# SLIDE 26: INVESTMENT
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, WARM_BG)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_GREEN)
add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Total 90-Day Investment: $10,000',
             font_size=28, bold=True, color=DARK_NAVY)
add_accent_bar(slide, 0.8, 1.15, 3.0, 0.04, ACCENT_GREEN)

stat_box(slide, 0.8, 1.6, 2.5, 1.5, '$2,000', 'Brand Design\n(Freelance)', ACCENT_TEAL)
stat_box(slide, 3.6, 1.6, 2.5, 1.5, '$2,000', 'Product\nPhotography', ACCENT_BLUE)
stat_box(slide, 6.4, 1.6, 2.5, 1.5, '$3,500', 'Shopify + POD\nSetup', ACCENT_ORANGE)
stat_box(slide, 9.2, 1.6, 2.5, 1.5, '$2,500', 'Workshop Pilot\n+ Misc', ACCENT_GREEN)

add_multi_text(slide, 0.8, 3.7, 11.5, 3.5, [
    ('Expected Return within 90 Days:', True, False),
    ('  $15-25K in incremental digital revenue from artist drops and Shopify optimization', False, False),
    ('  3 documented artist case studies (evidence for funder presentation)', False, False),
    ('  Validated (or invalidated) workshop hypothesis', False, False),
    ('  Competitive positioning vs GraffitiStreet, DirtyPilot, Etsy, StockX', False, False),
    ('  Partner gallery conversations started (Alberta St, Gray Loft, Side Rail)', False, False),
    ('', False, False),
    ('This $10K transforms Rossi from a business plan into a validated, evidence-backed funding proposal.', True, True),
    ('The evidence unlocks $185K in staged funding.', True, False),
], font_size=15, color=DARK_TEXT, line_spacing=1.15)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: Underwriter Audit \u2014 Investment Framework; Marketplace Research \u2014 Competitor Data]')

# ================================================================
# SLIDE 27: SUCCESS METRICS
# ================================================================
slide = two_column_slide(
    'Success Metrics: Outcomes, Not Outputs',
    'Customer Value Metrics',
    [
        ('Artist income from Rossi: $2,500 avg', False, False),
        ('Top-tier artist income: $15K+ (top 8)', False, False),
        ('Workshop satisfaction: 4.5/5', False, False),
        ('Repeat collector rate: 20%+', False, False),
        ('Artist retention rate: 85%+', False, False),
        ('', False, False),
        ('Competitive benchmark:', True, False),
        ('  GraffitiStreet takes 40-50%. Etsy 6.5%.', False, False),
        ('  Rossi at 50/50 must prove the', False, False),
        ('  development value justifies the split.', False, False),
    ],
    'Business Sustainability',
    [
        ('Total revenue: $745K', False, False),
        ('Online revenue: $55K', False, False),
        ('Cash balance: >$25K at all times', False, False),
        ('Revenue diversification: no stream >50%', False, False),
        ('Grant pipeline: $200K applied, $110K received', False, False),
        ('', False, False),
        ('90-Day Discovery Targets:', True, False),
        ('  Workshop: 8+ attendees', False, False),
        ('  First artist drop: $500+ in 72hrs', False, False),
        ('  3 case studies + 1 partner LOI', False, False),
    ],
    accent_color=ACCENT_GREEN
)
add_evidence(slide, 0.8, 7.0, 11, '[Evidence: BCG Assessment \u2014 Success Metrics; Marketplace Research \u2014 Commission Benchmarks]')

# ================================================================
# SLIDE 28: SECTION — RECOMMENDATION
# ================================================================
section_slide('Recommendation', '')

# ================================================================
# SLIDE 29: THE VERDICT
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, DARK_NAVY)

add_text_box(slide, 1.5, 0.6, 10, 0.8,
             'Recommendation to Rossi Leadership',
             font_size=30, bold=True, color=WHITE)
add_accent_bar(slide, 1.5, 1.5, 3.0, 0.06, ACCENT_TEAL)

# Verdict box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.8),
    Inches(9.3), Inches(1.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x22, 0x3A, 0x22)
shape.line.color.rgb = ACCENT_GREEN
shape.line.width = Pt(3)

tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'CONDITIONAL GO'
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = 'Invest 90 days and $10,000 in structured product discovery'
p2.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xAA)
p2.font.name = 'Calibri'
p2.alignment = PP_ALIGN.CENTER

add_multi_text(slide, 1.5, 3.6, 10, 3.5, [
    ('What Must Be True:', True, False),
    ('  1. Pull actual POS data within 2 weeks', False, False),
    ('  2. Run 1 pilot workshop within 6 weeks', False, False),
    ('  3. Document 3 real artist case studies', False, False),
    ('  4. Launch functioning e-commerce with monthly drops', False, False),
    ('  5. Resolve identity: cultural platform, not "nonprofit + brand + incubator"', False, False),
    ('  6. Moderate growth to 35-45% annually (nonprofit benchmark)', False, False),
    ('  7. Restructure partner economics before expansion', False, False),
    ('  8. Benchmark against GraffitiStreet, DirtyPilot, Alberta St, Side Rail', False, False),
], font_size=14, color=RGBColor(0xDD, 0xDD, 0xEE), line_spacing=1.15)

add_text_box(slide, 1.5, 6.8, 10, 0.4,
             'Underwriter Score: 82/100  |  BCG Verdict: Conditional Go  |  Palette System: Validated',
             font_size=11, italic=True, color=EVIDENCE_COLOR,
             alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 30: CLOSING
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, DARK_NAVY)

add_text_box(slide, 1.5, 2.0, 10, 1.0,
             'The gap is not vision.',
             font_size=32, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.8, 10, 1.0,
             'Rossi has exceptional vision.',
             font_size=24, italic=True, color=RGBColor(0xAA, 0xAA, 0xCC),
             alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, 5.5, 3.8, 2.3, 0.06, ACCENT_TEAL)

add_text_box(slide, 1.5, 4.2, 10, 1.0,
             'The gap is evidence.',
             font_size=32, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.0, 10, 1.0,
             'Close it in 90 days, and this becomes a funded product.',
             font_size=24, italic=True, color=ACCENT_TEAL,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 6.0, 10, 0.5,
             'V2 \u2014 Validated with Palette System  |  Para Signal Monitor: 6/6 anomalies resolved',
             font_size=11, color=RGBColor(0x66, 0x88, 0x66),
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 6.5, 10, 0.5,
             'Product Management Practice  |  February 2026  |  Confidential',
             font_size=11, color=RGBColor(0x66, 0x66, 0x88),
             alignment=PP_ALIGN.CENTER)

# ================================================================
# SAVE
# ================================================================
output_path = '/home/mical/fde/rossi-mission-project/ROSSI_PRODUCT_STRATEGY.pptx'
prs.save(output_path)
print(f'Slide deck saved to: {output_path}')
