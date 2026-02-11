#!/usr/bin/env python3
"""Generate AI Use Case & Metrics slide deck.

All slides: blank layout + background image + textboxes with run-level formatting.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# -- Paths --
OUTPUT = '/home/mical/fde/projects/agent-class/interview-prep/AGENTIC_SYSTEMS_CLASS.pptx'
BG_IMG = '/home/mical/fde/projects/agent-class/interview-prep/backround.png'

# -- Create blank widescreen presentation --
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

LAYOUT_BLANK = prs.slide_layouts[6]  # Blank layout

# -- Colors --
TITLE_BLUE = RGBColor(0x00, 0xB0, 0xF0)
ACCENT_BLUE = RGBColor(0x38, 0x96, 0xD2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MED_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT_TEAL = RGBColor(0x2A, 0x9D, 0x8F)
ACCENT_ORANGE = RGBColor(0xE7, 0x6F, 0x51)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_YELLOW = RGBColor(0xF3, 0x9C, 0x12)

FONT = 'Fira Sans'
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def new_slide(notes=None):
    """Create a blank slide with the background image."""
    slide = prs.slides.add_slide(LAYOUT_BLANK)
    slide.shapes.add_picture(BG_IMG, 0, 0, SLIDE_W, SLIDE_H)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def _set_run(p, text, font_size, bold, color, font_name=FONT):
    """Set text on a paragraph via a run with explicit formatting."""
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_title(slide, text, top=0.25, font_size=36, color=TITLE_BLUE,
              alignment=PP_ALIGN.LEFT, left=0.5, width=9.0):
    """Add a title textbox."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    _set_run(p, text, font_size, True, color)
    return txBox


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    """Add a simple text box with one line."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    _set_run(p, text, font_size, bold, color)
    return txBox


def add_body(slide, items, top=1.1, left=0.6, width=8.8, font_size=16,
             default_color=WHITE):
    """Add bulleted body content via runs.

    items: list of (text, bold, indent) or (text, bold, indent, color) tuples.
    """
    height = 5.625 - top - 0.3
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if len(item) == 4:
            text, bold, indent, color = item
        elif len(item) == 3:
            text, bold, indent = item
            color = default_color
        else:
            text, bold = item
            indent = 0
            color = default_color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.LEFT

        # Indent via XML margin
        if indent > 0:
            pPr = p._p.get_or_add_pPr()
            pPr.set(qn('a:marL'), str(int(Inches(0.35 * indent))))

        # Set text via run (this is what makes it visible in LibreOffice)
        _set_run(p, text, font_size, bold, color)

    return txBox


def add_multi_textbox(slide, left, top, width, height, lines, font_size=16,
                      color=WHITE, alignment=PP_ALIGN.LEFT, spacing=Pt(4)):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if len(line) == 3:
            text, bold, line_color = line
        else:
            text, bold = line
            line_color = color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = spacing
        _set_run(p, text, font_size, bold, line_color)
    return txBox


def content_slide(title_text, body_items, title_size=36, body_size=16, notes=None,
                  body_width=8.8):
    """Create a slide with title + bulleted body."""
    slide = new_slide(notes=notes)
    add_title(slide, title_text, font_size=title_size)
    add_body(slide, body_items, font_size=body_size, width=body_width)
    return slide


def section_slide(title_text, subtitle_text=None, notes=None):
    """Create a section divider slide."""
    slide = new_slide(notes=notes)
    add_textbox(slide, 0.5, 1.5, 9.0, 1.0, title_text,
                font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle_text:
        add_textbox(slide, 0.5, 2.6, 9.0, 0.6, subtitle_text,
                    font_size=20, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def add_rounded_box(slide, left, top, width, height, text, fill_color,
                    text_color=WHITE, font_size=12, bold=True):
    """Add a rounded rectangle with centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p, text, font_size, bold, text_color)
    return shape


def add_connector(slide, x1, y1, x2, y2, color=LIGHT_GRAY):
    """Add a line connector."""
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(2)
    return c


def draw_four_question_flowchart(slide):
    """Draw the 4-Question Method flowchart on the right side."""
    x, bw, bh, gap = 6.8, 2.7, 0.45, 0.12
    add_rounded_box(slide, x, 0.7, bw, 0.40, "4-QUESTION METHOD",
                    ACCENT_BLUE, WHITE, 10, True)
    questions = [("1. Pattern?", ACCENT_TEAL), ("2. Data?", ACCENT_TEAL),
                 ("3. Volume?", ACCENT_TEAL), ("4. Measurable?", ACCENT_TEAL)]
    y = 1.25
    for i, (text, color) in enumerate(questions):
        add_rounded_box(slide, x, y, bw, bh, text, color, WHITE, 10)
        if i > 0:
            add_connector(slide, x + bw/2, y - gap, x + bw/2, y, LIGHT_GRAY)
        y += bh + gap
    add_connector(slide, x + bw/2, y - gap, x + bw/2, y, LIGHT_GRAY)
    add_rounded_box(slide, x, y, bw, bh, "ALL YES -> PROCEED", ACCENT_GREEN, WHITE, 10, True)
    y += bh + gap
    add_connector(slide, x + bw/2, y - gap, x + bw/2, y, LIGHT_GRAY)
    add_rounded_box(slide, x, y, bw, bh, "ANY NO -> DIG DEEPER", ACCENT_ORANGE, WHITE, 10, True)


def draw_validation_tiers(slide):
    """Draw market validation tiers on the right side."""
    x, bw, bh, gap = 6.8, 2.7, 0.75, 0.1
    tiers = [("HIGH", "20+ co's, $500M+", ACCENT_GREEN),
             ("MEDIUM", "5-20 co's, $100-500M", ACCENT_YELLOW),
             ("LOW / GAP", "0-5 co's, <$50M", ACCENT_RED)]
    y = 0.8
    for title, sub, color in tiers:
        shape = add_rounded_box(slide, x, y, bw, bh, "", color, WHITE, 10)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        _set_run(p1, title, 11, True, WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _set_run(p2, sub, 9, False, WHITE)
        y += bh + gap
    add_textbox(slide, x, y + 0.05, bw, 0.3, "Market Validation Signal",
                font_size=9, bold=True, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def draw_metrics_timeline(slide):
    """Draw metrics timeline on the right side."""
    x, bw, bh, gap = 6.8, 2.7, 0.55, 0.08
    phases = [("0-3 MONTHS", "Accuracy, adoption", ACCENT_TEAL),
              ("3-6 MONTHS", "Business impact", ACCENT_BLUE),
              ("6+ MONTHS", "ROI, expansion", ACCENT_GREEN)]
    y = 0.8
    for i, (title, sub, color) in enumerate(phases):
        shape = add_rounded_box(slide, x, y, bw, bh, "", color, WHITE, 10)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        _set_run(p1, title, 11, True, WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _set_run(p2, sub, 9, False, WHITE)
        if i > 0:
            add_connector(slide, x + bw/2, y - gap, x + bw/2, y, LIGHT_GRAY)
        y += bh + gap


# ================================================================
# SLIDE 1: TITLE
# ================================================================
slide = new_slide(notes="Gagne P1: Gain Attention")
add_textbox(slide, 0.5, 1.2, 9.0, 1.0,
            "AI Use Case & Metrics",
            font_size=44, bold=True, color=TITLE_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.3, 9.0, 0.6,
            "Problem Framing  |  Success Metrics  |  MVP Scope",
            font_size=20, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 3.1, 9.0, 0.5,
            "Instructor: Mical Neill  |  Interview Prep / Live Class  |  Feb 2026",
            font_size=16, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2: INSTRUCTOR INTRODUCTION
# ================================================================
slide = new_slide()
add_title(slide, "Instructor Introduction")
add_textbox(slide, 0.5, 1.0, 9.0, 0.5,
            "Mical Neill", font_size=24, bold=True, color=WHITE)
add_multi_textbox(slide, 0.5, 1.6, 9.0, 3.5,
    lines=[
        ("7 years in AWS Knowledge Architecture", False),
        ("Currently building multi-agent systems for product strategy and research", False),
        ("", False),
        ("Background:", True, TITLE_BLUE),
        ("  Designed knowledge systems at scale (AWS)", False),
        ("  Built and shipped agentic AI applications in production", False),
        ("  Mapped 127 AI companies across 12 use cases for this class", False),
        ("", False),
        ("Teaching Philosophy:", True, TITLE_BLUE),
        ("  Learn patterns, not memorized answers", False),
        ("  Prepare for unseen problems with transferable frameworks", False),
    ],
    font_size=16, color=WHITE, spacing=Pt(3))


# ================================================================
# SLIDE 3: TODAY'S AGENDA
# ================================================================
content_slide(
    "Today's Agenda",
    [
        ("3 modules, one framework you can use in any AI interview:", True, 0, TITLE_BLUE),
        ("", False, 0),
        ("1.  Problem Framing and AI Fit", True, 0),
        ("How to determine if a problem is right for AI", False, 1),
        ("The 4-Question Method + market validation with real data", False, 1),
        ("", False, 0),
        ("2.  Success Metrics and Evaluation Readiness", True, 0),
        ("Time-phased metrics (0-3 mo, 3-6 mo, 6+ mo)", False, 1),
        ("Two live interview scenarios: mature vs. emerging market", False, 1),
        ("", False, 0),
        ("3.  MVP Scope and Prioritization", True, 0),
        ("Scoping principles and the 80/20 rule for AI MVPs", False, 1),
        ("Common mistakes and how to avoid them", False, 1),
    ],
    body_size=17,
    notes="Gagne P2: State the Learning Objective."
)


# ================================================================
# SECTION 1: PROBLEM FRAMING AND AI FIT
# ================================================================
section_slide(
    "Problem Framing and AI Fit",
    "How to determine if a problem is right for AI",
    notes="Gagne P1: Gain Attention."
)


# ================================================================
# SLIDE 5: THE INTERVIEW TRAP
# ================================================================
content_slide(
    "The Interview Trap",
    [
        ("Common interview questions:", True, 0, TITLE_BLUE),
        ('"How would you evaluate if [problem] is a good fit for AI?"', False, 1),
        ('"A customer wants to build [AI solution]. How do you assess viability?"', False, 1),
        ('"What metrics would you use to determine success?"', False, 1),
        ("", False, 0),
        ("The Trap:", True, 0, ACCENT_RED),
        ("Jumping to solutions without validating the problem", False, 1),
        ('"We\'ll use GPT-4 with RAG and fine-tune on their data..."', False, 1),
        ("", False, 0),
        ("What Interviewers Actually Want:", True, 0, ACCENT_GREEN),
        ("Structured thinking -- a repeatable framework", False, 1),
        ("Risk awareness -- not just optimism", False, 1),
        ("Metrics-first approach -- define success before building", False, 1),
    ],
    title_size=34, body_size=16,
    notes="Gagne P1: Gain Attention."
)


# ================================================================
# SLIDE 6: THE 4-QUESTION METHOD (with flowchart)
# ================================================================
slide = content_slide(
    "The 4-Question Method",
    [
        ("Run every AI use case through 4 questions:", True, 0, TITLE_BLUE),
        ("", False, 0),
        ("1.  Is there a pattern?", True, 0, ACCENT_TEAL),
        ("Can this be solved by recognizing patterns in data?", False, 1),
        ("", False, 0),
        ("2.  Is there data?", True, 0, ACCENT_TEAL),
        ("Enough quality data to train, fine-tune, or retrieve from?", False, 1),
        ("", False, 0),
        ("3.  Is it repetitive or high-volume?", True, 0, ACCENT_TEAL),
        ("Does this happen often enough to justify AI investment?", False, 1),
        ("", False, 0),
        ("4.  Can you measure success?", True, 0, ACCENT_TEAL),
        ("Clear metrics to know if AI is working?", False, 1),
        ("", False, 0),
        ("All 4 must be YES or CHECKABLE", True, 0, ACCENT_ORANGE),
    ],
    title_size=34, body_size=15, body_width=6.0,
    notes="Gagne P3: Present Information -- core framework."
)
draw_four_question_flowchart(slide)


# ================================================================
# SLIDE 7: Q1 & Q2 DEEP DIVE
# ================================================================
content_slide(
    "Pattern Recognition  +  Data Availability",
    [
        ("Good AI Fit -- Pattern:", True, 0, ACCENT_GREEN),
        ("Customer support tickets -- common questions recur", False, 1),
        ("Code completion -- syntax and context are predictable", False, 1),
        ("Document extraction -- invoices have structure", False, 1),
        ("", False, 0),
        ("Poor AI Fit -- No Pattern:", True, 0, ACCENT_RED),
        ("One-off strategic decisions, creative brand positioning, novel research", False, 1),
        ("", False, 0),
        ("Good AI Fit -- Data Available:", True, 0, ACCENT_GREEN),
        ("Historical conversations (thousands of examples)", False, 1),
        ("Public code repos (billions of lines)", False, 1),
        ("Existing documents (structured/semi-structured)", False, 1),
        ("", False, 0),
        ("Poor AI Fit -- No Data:", True, 0, ACCENT_RED),
        ("Brand new process, highly sensitive data, rare events", False, 1),
        ("", False, 0),
        ("Key insight: AI is pattern recognition, not magic", True, 0, ACCENT_ORANGE),
    ],
    title_size=30, body_size=14
)


# ================================================================
# SLIDE 8: Q3 & Q4 DEEP DIVE
# ================================================================
content_slide(
    "Volume & Repetition  +  Measurability",
    [
        ("Good AI Fit -- High Volume:", True, 0, ACCENT_GREEN),
        ("10,000+ support tickets/month", False, 1),
        ("Developers writing code daily", False, 1),
        ("1,000+ documents processed per week", False, 1),
        ("", False, 0),
        ("Poor AI Fit -- Low Volume:", True, 0, ACCENT_RED),
        ("Quarterly strategic planning (4x/year)", False, 1),
        ("Annual compliance audits (1x/year)", False, 1),
        ("", False, 0),
        ("Good Metrics -- Measurable:", True, 0, ACCENT_GREEN),
        ("Support: resolution time, CSAT, ticket deflection rate", False, 1),
        ("Code: acceptance rate, time saved, bugs introduced", False, 1),
        ("Documents: extraction accuracy, processing time", False, 1),
        ("", False, 0),
        ("Poor Metrics -- Vague:", True, 0, ACCENT_RED),
        ("\"Make our brand more innovative\" or \"improve morale\"", False, 1),
        ("", False, 0),
        ("Key insight: ROI requires volume. Define metrics before building.", True, 0, ACCENT_ORANGE),
    ],
    title_size=30, body_size=14
)


# ================================================================
# SLIDE 9: MARKET VALIDATION (with tier diagram)
# ================================================================
slide = content_slide(
    "Market Validation: Funded Companies = Signal",
    [
        ("127 companies mapped across 12 AI use cases", True, 0, TITLE_BLUE),
        ("", False, 0),
        ("HIGH -- 20+ companies, $500M+", True, 0, ACCENT_GREEN),
        ("Sales & GTM: 25+ co's, $2B+", False, 1),
        ("Customer Support: 20+ co's, $700M+", False, 1),
        ("Dev Productivity: 30+ co's, $500M+", False, 1),
        ("Risk: Differentiation in crowded market", False, 1, ACCENT_ORANGE),
        ("", False, 0),
        ("MEDIUM -- 5-20 companies, $100M-$500M", True, 0, ACCENT_YELLOW),
        ("Agentic Orchestration: 15 co's, $150M+", False, 1),
        ("AI Governance: 8-10 co's", False, 1),
        ("Risk: Category still being defined", False, 1, ACCENT_ORANGE),
        ("", False, 0),
        ("LOW / White Space -- 0-5 companies, <$50M", True, 0, ACCENT_RED),
        ("AI Decision Logs, Internal Demo Gen", False, 1),
        ("Risk: Market may not exist", False, 1, ACCENT_ORANGE),
    ],
    title_size=28, body_size=13, body_width=6.0,
    notes="Gagne P3: Present Information -- market data."
)
draw_validation_tiers(slide)


# ================================================================
# SLIDE 10: REAL COMPANY DATA
# ================================================================
content_slide(
    "Real Data: 127 Companies Across 12 Use Cases",
    [
        ("Highest Funded Use Cases:", True, 0, TITLE_BLUE),
        ("Sales & GTM: Gong $584M, Clari $576M, Outreach $489M", False, 1),
        ("Customer Support: Intercom $240M, Ada $190M, Kustomer $174M", False, 1),
        ("Document Processing: Hyperscience $218M, Instabase $132M", False, 1),
        ("RAG Infrastructure: Pinecone $138M, Weaviate $67M", False, 1),
        ("Knowledge Mgmt: Glean $360M, Notion AI $343M, Perplexity $100M", False, 1),
        ("", False, 0),
        ("Emerging / Agent-Native:", True, 0, ACCENT_TEAL),
        ("Orchestration: LangChain $35M, Fixie $17M, Dust $16M", False, 1),
        ("Workflow Automation: Zapier $1.3B, Make $38M, Bardeen $18M", False, 1),
        ("AI Evaluation: Patronus $17M, Braintrust $5M, Humanloop $5M", False, 1),
        ("", False, 0),
        ("New Unicorns (2023-2025):", True, 0, ACCENT_ORANGE),
        ("Sierra: $525M, $10B valuation in 2 years (customer support)", False, 1),
        ("Cursor: $900M Series C, ~$10B valuation (code generation)", False, 1),
        ("Harvey: $5B valuation (legal AI)", False, 1),
    ],
    title_size=28, body_size=13
)


# ================================================================
# SECTION 2: SUCCESS METRICS AND EVALUATION READINESS
# ================================================================
section_slide(
    "Success Metrics and Evaluation Readiness",
    "Time-phased metrics and real interview scenarios",
    notes="Gagne P3/P4: Present Information + Assess Performance."
)


# ================================================================
# SLIDE 12: TIME-PHASED METRICS (with timeline)
# ================================================================
slide = content_slide(
    "Time-Phased Success Metrics",
    [
        ("Metrics should be concrete and time-phased:", True, 0, TITLE_BLUE),
        ("", False, 0),
        ("Early (0-3 months):", True, 0, ACCENT_TEAL),
        ("Accuracy on core task (e.g., 85% on top 10 question types)", False, 1),
        ("User adoption rate, error rate baseline", False, 1),
        ("", False, 0),
        ("Medium (3-6 months):", True, 0, ACCENT_TEAL),
        ("Business impact (e.g., 30% ticket deflection, 50% time saved)", False, 1),
        ("Active users at scale, cost per interaction", False, 1),
        ("", False, 0),
        ("Long-term (6+ months):", True, 0, ACCENT_TEAL),
        ("Sustained quality (CSAT maintained), cost reduction", False, 1),
        ("Expansion to adjacent use cases", False, 1),
        ("", False, 0),
        ("Evaluation readiness: can you measure TODAY?", True, 0, ACCENT_ORANGE),
        ("If no baseline data exists, that's your first milestone", False, 1),
    ],
    title_size=32, body_size=15, body_width=6.0
)
draw_metrics_timeline(slide)


# ================================================================
# SLIDE 13: SCENARIO 1 -- MATURE MARKET
# ================================================================
content_slide(
    "Scenario 1: Mature Market",
    [
        ("Interview question:", True, 0, TITLE_BLUE),
        ('"A customer wants to build an AI-powered customer support chatbot."', False, 1),
        ('"How would you evaluate if this is a good idea?"', False, 1),
        ("", False, 0),
        ("4-Question Method:", True, 0, ACCENT_TEAL),
        ("Pattern? YES -- common questions, resolution paths recur", False, 1),
        ("Data? CHECK -- need 6+ months of historical tickets", False, 1),
        ("Repetitive? CHECK -- need 1,000+ tickets/month", False, 1),
        ("Measurable? YES -- deflection rate, CSAT, resolution time", False, 1),
        ("", False, 0),
        ("Market Validation: 20+ companies, $700M+ -- HIGH", True, 0, ACCENT_GREEN),
        ("Intercom $240M, Ada $190M, Kustomer $174M, Forethought $92M", False, 1),
        ("Challenge is differentiation, not viability", False, 1),
        ("", False, 0),
        ("Metrics: 85% accuracy (0-3mo) -> 30% deflection (3-6mo) -> CSAT maintained (6+mo)", True, 0, ACCENT_ORANGE),
    ],
    title_size=30, body_size=14,
    notes="Gagne P4: Assess Performance."
)


# ================================================================
# SLIDE 14: SCENARIO 2 -- EMERGING CATEGORY
# ================================================================
content_slide(
    "Scenario 2: Emerging Category",
    [
        ("Interview question:", True, 0, TITLE_BLUE),
        ('"A customer wants to build an agentic workflow automation platform."', False, 1),
        ('"How would you evaluate this?"', False, 1),
        ("", False, 0),
        ("4-Question Method:", True, 0, ACCENT_TEAL),
        ("Pattern? MAYBE -- workflows have patterns, but complex", False, 1),
        ("Data? PARTIAL -- individual tools yes, cross-tool less", False, 1),
        ("Repetitive? YES -- if targeting common workflows", False, 1),
        ("Measurable? YES -- time saved, error rate", False, 1),
        ("", False, 0),
        ("Market Validation: 15 co's, $150M+ -- EMERGING", True, 0, ACCENT_YELLOW),
        ("LangChain $35M, Fixie $17M, Dust $16M -- all founded 2022-23", False, 1),
        ("", False, 0),
        ("Metrics: 70% completion (0-3mo) -> 50% time saved (3-6mo)", True, 0, ACCENT_ORANGE),
        ("", False, 0),
        ("Key difference:", True, 0, WHITE),
        ("Mature market = differentiation risk", False, 1, ACCENT_GREEN),
        ("Emerging market = category + reliability risk", False, 1, ACCENT_YELLOW),
    ],
    title_size=30, body_size=13,
    notes="Gagne P4: Assess Performance -- contrast."
)


# ================================================================
# SECTION 3: MVP SCOPE AND PRIORITIZATION
# ================================================================
section_slide(
    "MVP Scope and Prioritization",
    "Scoping discipline separates good answers from great ones",
    notes="Gagne P5: Provide Feedback."
)


# ================================================================
# SLIDE 16: MVP SCOPING PRINCIPLES
# ================================================================
content_slide(
    "MVP Scoping Principles",
    [
        ("The 80/20 rule for AI MVPs:", True, 0, TITLE_BLUE),
        ("Solve the top 10 use cases, not all of them", False, 1),
        ("", False, 0),
        ("Scoping rules:", True, 0, ACCENT_TEAL),
        ("ONE workflow, not \"automate everything\"", False, 1),
        ("3-5 integrations, not the entire tool ecosystem", False, 1),
        ("Human-in-the-loop for edge cases", False, 1),
        ("", False, 0),
        ("Prioritization criteria:", True, 0, ACCENT_TEAL),
        ("Volume -- how often does this task happen?", False, 1),
        ("Data readiness -- is training/retrieval data available?", False, 1),
        ("Measurability -- can we prove success in 90 days?", False, 1),
        ("Risk tolerance -- what happens when AI is wrong?", False, 1),
        ("", False, 0),
        ("Red flag: if you can't scope the MVP to one sentence, it's too broad", True, 0, ACCENT_RED),
    ],
    title_size=32, body_size=15
)


# ================================================================
# SLIDE 17: COMMON MISTAKES
# ================================================================
content_slide(
    "Common Mistakes in AI Use Case Interviews",
    [
        ("1. Jumping to solutions", True, 0, ACCENT_RED),
        ("BAD: \"We'll use GPT-4 with RAG...\"", False, 1),
        ("GOOD: \"Let me validate AI suitability first...\"", False, 1, ACCENT_GREEN),
        ("", False, 0),
        ("2. Ignoring market validation", True, 0, ACCENT_RED),
        ("BAD: \"Great idea, let's build!\"", False, 1),
        ("GOOD: \"I see 20+ companies with $700M funding -- differentiation is key\"", False, 1, ACCENT_GREEN),
        ("", False, 0),
        ("3. Vague metrics", True, 0, ACCENT_RED),
        ("BAD: \"Measure user satisfaction\"", False, 1),
        ("GOOD: \"85% accuracy on top 10, 30% deflection rate by month 6\"", False, 1, ACCENT_GREEN),
        ("", False, 0),
        ("4. Ignoring risks", True, 0, ACCENT_RED),
        ("BAD: \"AI can solve this easily!\"", False, 1),
        ("GOOD: \"Risk is differentiation in a crowded market\"", False, 1, ACCENT_GREEN),
        ("", False, 0),
        ("5. Over-scoping the MVP", True, 0, ACCENT_RED),
        ("BAD: \"Automate all support\"", False, 1),
        ("GOOD: \"MVP: Top 10 question types\"", False, 1, ACCENT_GREEN),
    ],
    title_size=28, body_size=14,
    notes="Gagne P5: Provide Feedback -- anti-patterns."
)


# ================================================================
# SLIDE 18: KEY TAKEAWAYS
# ================================================================
content_slide(
    "Key Takeaways",
    [
        ("1.  Use the 4-Question Method every time", True, 0),
        ("Pattern -> Data -> Volume -> Measurable -- all must pass", False, 1),
        ("", False, 0),
        ("2.  Market validation = funded companies", True, 0),
        ("High validation = differentiation risk. Low = existence risk.", False, 1),
        ("", False, 0),
        ("3.  Time-phase your metrics", True, 0),
        ("0-3 mo: accuracy. 3-6 mo: impact. 6+ mo: sustained ROI.", False, 1),
        ("", False, 0),
        ("4.  Scope aggressively for MVP", True, 0),
        ("One workflow. Measurable in 90 days. Human-in-the-loop.", False, 1),
        ("", False, 0),
        ("5.  Name risks explicitly", True, 0),
        ("Mature market != emerging market -- different risks, different recs.", False, 1),
        ("", False, 0),
        ("What interviewers look for:", True, 0, TITLE_BLUE),
        ("PM: customer empathy + prioritization", False, 1),
        ("TPM: risk identification + execution planning", False, 1),
        ("EM: technical depth + reliability awareness", False, 1),
    ],
    title_size=34, body_size=14,
    notes="Gagne P5: Provide Feedback -- 5 actionable takeaways."
)


# ================================================================
# SLIDE 19: LIVE DEMO
# ================================================================
slide = new_slide(notes="Gagne P4/P5: Live application.")
add_title(slide, "Live Demo: AI Use Case Evaluator", font_size=32)
add_multi_textbox(slide, 0.5, 1.0, 9.0, 4.2,
    lines=[
        ("You just learned the framework. Now watch an agent apply it.", True),
        ("", False),
        ("One prompt. Any use case. Full evaluation in ~30 seconds.", False),
        ("", False),
        ("What it produces:", True, TITLE_BLUE),
        ("  4-Question Method analysis (pattern, data, volume, measurable)", False),
        ("  Market validation signal (company count + funding tier)", False),
        ("  MVP scope (one sentence + constraints)", False),
        ("  Time-phased metrics (0-3mo, 3-6mo, 6+mo)", False),
        ("  Risk assessment + recommendation (go / conditional / no-go)", False),
        ("", False),
        ("Audience: suggest a use case.", True, ACCENT_ORANGE),
        ("", False),
        ("  \"AI-powered resume screening for recruiters\"", False),
        ("  \"AI that detects fraudulent insurance claims\"", False),
        ("  \"AI copilot for financial advisors\"", False),
    ],
    font_size=16, color=WHITE, spacing=Pt(3))


# ================================================================
# SLIDE 20: QUESTIONS
# ================================================================
slide = new_slide()
add_textbox(slide, 0.5, 1.5, 9.0, 1.0, "Questions?",
            font_size=44, bold=True, color=TITLE_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 3.0, 9.0, 0.6,
            "Framework works for any AI use case -- practice with real scenarios",
            font_size=18, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 4.0, 9.0, 0.5,
            "Company-RIU Mapping Library: 127 AI companies across 12 use cases",
            font_size=14, bold=False, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ================================================================
# SAVE
# ================================================================
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Slide deck saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
